import os
import traceback
from fastapi import FastAPI, UploadFile, File, Form, HTTPException
from fastapi.staticfiles import StaticFiles
from fastapi.responses import HTMLResponse, FileResponse
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel, Field
from typing import Optional, List

from .agents.google_adk import SessionState
from .agents.coordinator import CoordinatorAgent
from .agents.planner import StudyPlannerAgent
from .agents.quiz import QuizAgent
from .agents.notes import NotesAgent
from .agents.progress import ProgressAgent
from .utils.security import (
    sanitize_input_text, 
    validate_uploaded_file, 
    UPLOAD_DIR, 
    ensure_sandbox_dirs, 
    BASE_WORKSPACE,
    detect_prompt_injection,
    validate_topic,
    sanitize_api_error,
    is_safe_path
)
from .utils.db import (
    load_db, 
    save_db, 
    get_user_profile, 
    save_user_profile, 
    get_study_plans, 
    add_study_plan, 
    delete_study_plan, 
    get_quiz_scores, 
    add_quiz_score, 
    get_uploaded_notes, 
    add_uploaded_note
)

# Ensure sandbox folders exist
ensure_sandbox_dirs()

# Initialize FastAPI
app = FastAPI(
    title="StudyPilot AI API Server",
    description="Offline-first multi-agent AI study concierge platform backend",
    version="1.0.0"
)

# Enable CORS for local development flexibility
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Instantiate Agents
coordinator_agent = CoordinatorAgent()
planner_agent = StudyPlannerAgent()
quiz_agent = QuizAgent()
notes_agent = NotesAgent()
progress_agent = ProgressAgent()

# --- Request/Response Schemas ---

class ChatRequest(BaseModel):
    query: str
    notes_text: Optional[str] = None
    notes_title: Optional[str] = None

class ChatResponse(BaseModel):
    response: str

class PlanCreateRequest(BaseModel):
    topic: str
    duration_weeks: int = Field(default=4, ge=1, le=12)
    daily_hours: float = Field(default=2.0, ge=0.5, le=8.0)
    skill_level: str = Field(default="Beginner")

class QuizCreateRequest(BaseModel):
    topic: str
    num_questions: int = Field(default=5, ge=1, le=10)
    notes_text: Optional[str] = None

class QuizScoreRequest(BaseModel):
    topic: str
    total_questions: int
    correct_answers: int
    score_percentage: float

class ProfileUpdateRequest(BaseModel):
    name: Optional[str] = None
    level: Optional[str] = None
    weekly_hours_goal: Optional[int] = None

# --- API Endpoints ---

@app.get("/api/profile")
async def get_profile():
    try:
        return get_user_profile()
    except Exception as e:
        print(f"Profile load error: {str(e)}\n{traceback.format_exc()}")
        raise HTTPException(status_code=500, detail="Failed to load student profile.")

@app.post("/api/profile")
async def update_profile(req: ProfileUpdateRequest):
    try:
        update_data = {}
        if req.name is not None:
            name_clean = sanitize_input_text(req.name)
            if not name_clean or len(name_clean) > 100:
                raise HTTPException(status_code=400, detail="Student name must be between 1 and 100 characters.")
            if not re.match(r"^[a-zA-Z0-9\s._\-]+$", name_clean):
                raise HTTPException(status_code=400, detail="Student name contains invalid characters.")
            update_data["name"] = name_clean
            
        if req.level is not None:
            level_clean = sanitize_input_text(req.level)
            if level_clean not in ["Beginner", "Intermediate", "Advanced"]:
                raise HTTPException(status_code=400, detail="Level must be Beginner, Intermediate, or Advanced.")
            update_data["level"] = level_clean
            
        if req.weekly_hours_goal is not None:
            update_data["weekly_hours_goal"] = max(1, min(100, int(req.weekly_hours_goal)))
            
        return save_user_profile(update_data)
    except HTTPException:
        raise
    except Exception as e:
        print(f"Profile update error: {str(e)}\n{traceback.format_exc()}")
        raise HTTPException(status_code=500, detail="Failed to update student profile.")

@app.post("/api/chat", response_model=ChatResponse)
async def chat_interaction(req: ChatRequest):
    try:
        # Validate query presence and length
        if not req.query or not req.query.strip():
            raise HTTPException(status_code=400, detail="Chat query cannot be empty.")
            
        if len(req.query) > 5000:
            raise HTTPException(status_code=400, detail="Chat query is too long (max 5000 characters).")
            
        # Prompt injection check
        is_injected, err_msg = detect_prompt_injection(req.query)
        if is_injected:
            raise HTTPException(status_code=400, detail=err_msg)
            
        # Sanitize chat input
        clean_query = sanitize_input_text(req.query)
        if not clean_query:
            raise HTTPException(status_code=400, detail="Query cannot contain solely HTML tags.")

        # Create session state and set context if notes are loaded
        session = SessionState()
        if req.notes_text:
            if len(req.notes_text) > 1024 * 1024:
                raise HTTPException(status_code=400, detail="Attached study notes exceed 1MB limit.")
            session.set("notes_text", req.notes_text)
            session.set("notes_title", sanitize_input_text(req.notes_title or "Uploaded File"))
            
        response = coordinator_agent.process_query(clean_query, session)
        return ChatResponse(response=response)
    except HTTPException:
        raise
    except Exception as e:
        error_clean = sanitize_api_error(str(e))
        print(f"Chat error: {error_clean}\n{traceback.format_exc()}")
        raise HTTPException(status_code=500, detail=f"Assistant error: {error_clean}")

@app.get("/api/planner")
async def list_plans():
    try:
        return get_study_plans()
    except Exception as e:
        print(f"List plans error: {str(e)}\n{traceback.format_exc()}")
        raise HTTPException(status_code=500, detail="Failed to load study plans.")

@app.post("/api/planner")
async def create_plan(req: PlanCreateRequest):
    try:
        print(f"DEBUG planner request: topic={req.topic!r}, weeks={req.duration_weeks!r}, hours={req.daily_hours!r}, level={req.skill_level!r}")
        
        # Extract core topic if it is a descriptive query (using coordinator_agent's extraction helper)
        extracted_topic = coordinator_agent._extract_topic(req.topic, req.topic, req.topic)
        print(f"DEBUG planner request: extracted_topic={extracted_topic!r}")
        
        # Validate topic
        is_valid_topic, err_msg_topic = validate_topic(extracted_topic)
        if not is_valid_topic:
            print(f"DEBUG planner validation failed on topic: {err_msg_topic}")
            raise HTTPException(status_code=400, detail=err_msg_topic)
            
        topic_san = sanitize_input_text(extracted_topic)
        level_san = sanitize_input_text(req.skill_level)

        
        if level_san not in ["Beginner", "Intermediate", "Advanced"]:
            print(f"DEBUG planner validation failed on level: {level_san}")
            raise HTTPException(status_code=400, detail="Skill level must be Beginner, Intermediate, or Advanced.")
            
        if not (1 <= req.duration_weeks <= 12):
            print(f"DEBUG planner validation failed on weeks: {req.duration_weeks}")
            raise HTTPException(status_code=400, detail="Duration must be between 1 and 12 weeks.")
            
        if not (0.5 <= req.daily_hours <= 8.0):
            print(f"DEBUG planner validation failed on hours: {req.daily_hours}")
            raise HTTPException(status_code=400, detail="Daily commitment must be between 0.5 and 8.0 hours.")
            
        # Execute Planner Agent
        plan_md = planner_agent.generate_plan(topic_san, req.duration_weeks, req.daily_hours, level_san)
        
        # Save to local database
        plan_entry = {
            "topic": topic_san,
            "duration_weeks": req.duration_weeks,
            "daily_hours": req.daily_hours,
            "skill_level": level_san,
            "plan_markdown": plan_md
        }
        add_study_plan(plan_entry)
        
        return plan_entry
    except HTTPException as he:
        print(f"DEBUG planner HTTPException: status={he.status_code}, detail={he.detail}")
        raise
    except Exception as e:
        error_clean = sanitize_api_error(str(e))
        print(f"Planner error: {error_clean}\n{traceback.format_exc()}")
        raise HTTPException(status_code=500, detail="Failed to generate study plan.")


@app.delete("/api/planner/{plan_id}")
async def remove_plan(plan_id: int):
    try:
        try:
            plan_id = int(plan_id)
        except ValueError:
            raise HTTPException(status_code=400, detail="Invalid plan ID.")
            
        success = delete_study_plan(plan_id)
        if not success:
            raise HTTPException(status_code=404, detail="Study plan not found.")
        return {"status": "success", "message": f"Plan {plan_id} deleted successfully."}
    except HTTPException:
        raise
    except Exception as e:
        print(f"Delete plan error: {str(e)}\n{traceback.format_exc()}")
        raise HTTPException(status_code=500, detail="Failed to delete study plan.")

@app.post("/api/quiz")
async def create_quiz(req: QuizCreateRequest):
    try:
        # Extract core topic if it is a descriptive query (using coordinator_agent's extraction helper)
        extracted_topic = coordinator_agent._extract_topic(req.topic, req.topic, req.topic)
        
        is_valid_topic, err_msg_topic = validate_topic(extracted_topic)
        if not is_valid_topic:
            raise HTTPException(status_code=400, detail=err_msg_topic)
            
        topic_san = sanitize_input_text(extracted_topic)
        
        if not (1 <= req.num_questions <= 10):
            raise HTTPException(status_code=400, detail="Quiz size must be between 1 and 10 questions.")
            
        if req.notes_text and len(req.notes_text) > 1024 * 1024:
            raise HTTPException(status_code=400, detail="Quiz context source notes exceed 1MB limit.")
            
        questions = quiz_agent.generate_quiz(topic_san, req.num_questions, req.notes_text)
        return {"topic": topic_san, "questions": questions}
    except HTTPException:
        raise
    except Exception as e:
        error_clean = sanitize_api_error(str(e))
        print(f"Quiz error: {error_clean}\n{traceback.format_exc()}")
        raise HTTPException(status_code=500, detail="Failed to generate practice quiz.")

@app.post("/api/quiz/score")
async def record_score(req: QuizScoreRequest):
    try:
        # Extract core topic if it is a descriptive query (using coordinator_agent's extraction helper)
        extracted_topic = coordinator_agent._extract_topic(req.topic, req.topic, req.topic)
        
        is_valid_topic, err_msg_topic = validate_topic(extracted_topic)
        if not is_valid_topic:
            raise HTTPException(status_code=400, detail=err_msg_topic)
            
        topic_san = sanitize_input_text(extracted_topic)
        
        if req.total_questions <= 0 or req.total_questions > 100:
            raise HTTPException(status_code=400, detail="Total questions must be between 1 and 100.")
            
        if req.correct_answers < 0 or req.correct_answers > req.total_questions:
            raise HTTPException(status_code=400, detail="Correct answers count is out of bounds.")
            
        # Verify score percentage calculation consistency
        expected_percentage = round((req.correct_answers / req.total_questions) * 100, 2)
        if abs(req.score_percentage - expected_percentage) > 1.0:
            raise HTTPException(status_code=400, detail="Incorrect score percentage provided.")
            
        score_entry = {
            "topic": topic_san,
            "total_questions": req.total_questions,
            "correct_answers": req.correct_answers,
            "score_percentage": req.score_percentage
        }
        add_quiz_score(score_entry)
        
        return progress_agent.analyze_progress()
    except HTTPException:
        raise
    except Exception as e:
        print(f"Record score error: {str(e)}\n{traceback.format_exc()}")
        raise HTTPException(status_code=500, detail="Failed to record quiz score.")

def extract_text_from_bytes(filename: str, content: bytes) -> str:
    _, ext = os.path.splitext(filename.lower())
    if ext in [".txt", ".md"]:
        return content.decode("utf-8", errors="ignore")
    elif ext == ".pdf":
        import pypdf
        import io
        pdf_file = io.BytesIO(content)
        reader = pypdf.PdfReader(pdf_file)
        text = []
        for page in reader.pages:
            t = page.extract_text()
            if t:
                text.append(t)
        return "\n".join(text)
    elif ext == ".docx":
        import docx
        import io
        doc_file = io.BytesIO(content)
        doc = docx.Document(doc_file)
        text = [p.text for p in doc.paragraphs]
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    text.append(cell.text)
        return "\n".join(text)
    elif ext == ".pptx":
        import pptx
        import io
        ppt_file = io.BytesIO(content)
        prs = pptx.Presentation(ppt_file)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text.append(shape.text)
        return "\n".join(text)
    elif ext == ".ppt":
        raise HTTPException(
            status_code=400, 
            detail="Format .ppt is currently unsupported. Please convert your PowerPoint presentation to .pptx format and try again."
        )
    else:
        raise HTTPException(status_code=400, detail=f"Unsupported file format: {ext}")

@app.post("/api/notes/upload")
async def upload_study_notes(file: UploadFile = File(...)):
    try:
        filename_san = os.path.basename(file.filename)
        file_path = os.path.join(UPLOAD_DIR, filename_san)
        
        # Verify safe path resolution
        if not is_safe_path(file_path, UPLOAD_DIR):
            raise HTTPException(status_code=400, detail="Invalid characters in file name (path traversal attempted).")
            
        # Validate MIME type
        allowed_mimes = {
            "text/plain",
            "text/markdown",
            "application/pdf",
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "application/vnd.ms-powerpoint",
            "application/octet-stream"
        }
        if file.content_type and file.content_type not in allowed_mimes:
            raise HTTPException(status_code=400, detail="MIME type validation failed. The uploaded file format is not supported.")
            
        content = await file.read()
        
        # Run security validations on uploaded note
        is_valid, err_msg = validate_uploaded_file(filename_san, content)
        if not is_valid:
            raise HTTPException(status_code=400, detail=err_msg)
            
        # Write file in sandbox uploads directory
        try:
            with open(file_path, "wb") as f:
                f.write(content)
        except Exception as e:
            print(f"File save exception: {str(e)}")
            raise HTTPException(status_code=500, detail="Failed to save uploaded file.")
            
        # Extract and parse text from file
        try:
            text_content = extract_text_from_bytes(filename_san, content)
        except HTTPException:
            raise
        except Exception as e:
            print(f"Error parsing file {filename_san}: {e}")
            raise HTTPException(status_code=400, detail="Failed to parse document text. Make sure the file is not corrupted.")
            
        summary_result = notes_agent.summarize_notes(text_content, filename_san)
        
        # Save note metadata in DB
        add_uploaded_note({
            "filename": filename_san,
            "word_count": len(text_content.split()),
            "keywords": summary_result["keywords"],
            "summary_markdown": summary_result["summary_markdown"]
        })
        
        return {
            "filename": filename_san,
            "summary_markdown": summary_result["summary_markdown"],
            "keywords": summary_result["keywords"],
            "full_content": text_content
        }
    except HTTPException:
        raise
    except Exception as e:
        print(f"Notes upload error: {str(e)}\n{traceback.format_exc()}")
        raise HTTPException(status_code=500, detail="Failed to upload and analyze notes.")

@app.get("/api/notes")
async def list_uploaded_notes():
    try:
        notes = get_uploaded_notes()
        notes_list = []
        for note in notes:
            filename = note.get("filename")
            file_path = os.path.join(UPLOAD_DIR, filename)
            
            # Read full content from sandbox directory if file exists
            full_content = ""
            if os.path.exists(file_path) and is_safe_path(file_path, UPLOAD_DIR):
                try:
                    with open(file_path, "rb") as f:
                        file_bytes = f.read()
                    full_content = extract_text_from_bytes(filename, file_bytes)
                except Exception as e:
                    print(f"Error reading note file {filename}: {e}")
            
            notes_list.append({
                "filename": filename,
                "summary_markdown": note.get("summary_markdown", ""),
                "keywords": note.get("keywords", []),
                "full_content": full_content
            })
        return notes_list
    except Exception as e:
        print(f"List notes error: {str(e)}\n{traceback.format_exc()}")
        raise HTTPException(status_code=500, detail="Failed to load study notes.")


@app.get("/api/progress")
async def get_progress():
    try:
        return progress_agent.analyze_progress()
    except Exception as e:
        print(f"Progress error: {str(e)}\n{traceback.format_exc()}")
        raise HTTPException(status_code=500, detail="Failed to load progress details.")

# --- Serve Static Frontend App ---

# Check if frontend directory exists in workspace
FRONTEND_DIR = os.path.join(BASE_WORKSPACE, "frontend")

@app.get("/")
async def serve_index():
    index_path = os.path.join(FRONTEND_DIR, "index.html")
    if os.path.exists(index_path):
        return FileResponse(index_path)
    return HTMLResponse("<h1>StudyPilot AI Dashboard</h1><p>Frontend code files not found yet. Please create index.html in the frontend folder.</p>")

# Mount other static assets if frontend exists
if os.path.exists(FRONTEND_DIR):
    app.mount("/", StaticFiles(directory=FRONTEND_DIR, html=True), name="static")

if __name__ == "__main__":
    import uvicorn
    # Run server on port 8000
    uvicorn.run("backend.app:app", host="127.0.0.1", port=8000, reload=True)
